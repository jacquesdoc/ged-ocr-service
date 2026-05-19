from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import easyocr
import tempfile
import os
from PIL import Image
import fitz
from docx import Document as DocxDocument
from pptx import Presentation

app = FastAPI(title="OCR Service - GED IVOPREST", version="2.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

print("Chargement du modele OCR...")
reader = easyocr.Reader(['fr', 'en'], gpu=False)
print("Modele OCR charge !")

@app.get("/")
def root():
    return {
        "service": "OCR Service",
        "version": "2.0.0",
        "status":  "running",
        "supported_formats": ["jpg", "jpeg", "png", "bmp", "tiff", "webp", "pdf", "docx", "pptx"]
    }

@app.get("/health")
def health():
    return {"status": "ok"}

def extract_pdf(file_path: str) -> dict:
    try:
        doc = fitz.open(file_path)
        text = ""
        page_count = len(doc)
        for page_num in range(page_count):
            page = doc[page_num]
            page_text = page.get_text()
            if page_text.strip():
                text += f"\n--- Page {page_num + 1} ---\n{page_text}"
            else:
                pix = page.get_pixmap(dpi=200)
                img_path = file_path + f"_page_{page_num}.png"
                pix.save(img_path)
                results = reader.readtext(img_path, detail=1)
                os.unlink(img_path)
                page_text = " ".join(t for (_, t, c) in results if c > 0.3)
                if page_text.strip():
                    text += f"\n--- Page {page_num + 1} (OCR) ---\n{page_text}"
        doc.close()
        text = text.strip()
        return {
            "success": True, "type": "pdf", "text": text,
            "page_count": page_count,
            "word_count": len(text.split()) if text else 0,
            "char_count": len(text), "confidence": 95 if text else 0,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur PDF : {str(e)}")

def extract_docx(file_path: str) -> dict:
    try:
        doc = DocxDocument(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        text = "\n".join(paragraphs)
        return {
            "success": True, "type": "docx", "text": text,
            "paragraph_count": len(paragraphs),
            "word_count": len(text.split()) if text else 0,
            "char_count": len(text), "confidence": 98,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur DOCX : {str(e)}")

def extract_pptx(file_path: str) -> dict:
    try:
        prs = Presentation(file_path)
        slides_text = []
        slide_count = len(prs.slides)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_content.append(cell.text.strip())
            if slide_content:
                slides_text.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_content))
        text = "\n\n".join(slides_text)
        return {
            "success": True, "type": "pptx", "text": text,
            "slide_count": slide_count,
            "word_count": len(text.split()) if text else 0,
            "char_count": len(text), "confidence": 98,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur PPTX : {str(e)}")

def extract_image(file_path: str, filename: str) -> dict:
    try:
        results = reader.readtext(file_path, detail=1)
        words = []
        for (bbox, text, confidence) in results:
            if confidence > 0.3:
                words.append({"text": text, "confidence": round(confidence * 100, 1)})
        text = " ".join(w["text"] for w in words).strip()
        avg_confidence = round(
            sum(w["confidence"] for w in words) / len(words), 1
        ) if words else 0
        return {
            "success": True, "type": "image", "text": text,
            "word_count": len(text.split()) if text else 0,
            "char_count": len(text), "confidence": avg_confidence,
            "words_detected": len(words),
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur image : {str(e)}")

@app.post("/ocr")
async def extract_text(file: UploadFile = File(...)):
    filename     = file.filename or "document"
    content_type = file.content_type or ""
    extension    = os.path.splitext(filename)[1].lower()

    image_types = ["image/jpeg", "image/jpg", "image/png",
                   "image/bmp", "image/tiff", "image/webp"]
    image_exts  = [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".webp"]
    pdf_types   = ["application/pdf"]
    docx_types  = [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword"
    ]
    pptx_types  = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint"
    ]

    contents = await file.read()

    with tempfile.NamedTemporaryFile(delete=False, suffix=extension) as tmp:
        tmp.write(contents)
        tmp_path = tmp.name

    try:
        if content_type in image_types or extension in image_exts:
            result = extract_image(tmp_path, filename)
        elif content_type in pdf_types or extension == ".pdf":
            result = extract_pdf(tmp_path)
        elif content_type in docx_types or extension in [".docx", ".doc"]:
            result = extract_docx(tmp_path)
        elif content_type in pptx_types or extension in [".pptx", ".ppt"]:
            result = extract_pptx(tmp_path)
        else:
            raise HTTPException(
                status_code=400,
                detail=f"Format non supporte : {content_type} ({extension})"
            )
        result["filename"] = filename
        return result
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)